#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PowerPoint Generator Module
===========================

This module creates professional PowerPoint presentations using python-pptx.
It provides functions for creating presentations with title slides, summary
slides, and chart slides.

Features:
    - Custom branded presentations
    - Title slides with subtitle and author
    - Summary slides with key metrics
    - Chart slides with images and descriptions

Usage:
    from ppt_generator import create_presentation, add_chart_slide, save_presentation
    
    prs = create_presentation("Report Title", "Subtitle", "Author")
    add_chart_slide(prs, "Chart Title", "chart.png", "Description")
    save_presentation(prs, "output.pptx")

Author: Data Analytics Team
License: MIT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path
from datetime import datetime


# =============================================================================
# Configuration
# =============================================================================

# Output directory for presentations
OUTPUT_DIR = Path("output/presentations")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# Brand colors
COLOR_PRIMARY = RGBColor(0x2E, 0x75, 0xB6)    # Blue
COLOR_SECONDARY = RGBColor(0x66, 0x66, 0x66)  # Gray
COLOR_LIGHT = RGBColor(0x99, 0x99, 0x99)      # Light gray


def create_presentation(
    title: str,
    subtitle: str = None,
    author: str = "Auto Report Generator"
) -> Presentation:
    """
    Create a new presentation with a branded title slide.
    
    Initializes a widescreen (16:9) presentation and adds a title slide
    with the specified title, subtitle, and author information.
    
    Args:
        title (str): Main presentation title.
        subtitle (str, optional): Subtitle text. Defaults to None.
        author (str, optional): Author/company name. Defaults to "Auto Report Generator".
        
    Returns:
        Presentation: python-pptx Presentation object ready for additional slides.
        
    Example:
        >>> prs = create_presentation("Q4 Sales Report", "2024 Performance", "Sales Team")
    """
    prs = Presentation()
    
    # Set widescreen dimensions (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Add title slide using blank layout for full control
    title_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(title_slide_layout)
    
    # --- Main Title ---
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRIMARY
    title_para.alignment = PP_ALIGN.CENTER
    
    # --- Subtitle ---
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(0.75))
        sub_frame = sub_box.text_frame
        sub_para = sub_frame.paragraphs[0]
        sub_para.text = subtitle
        sub_para.font.size = Pt(24)
        sub_para.font.color.rgb = COLOR_SECONDARY
        sub_para.alignment = PP_ALIGN.CENTER
    
    # --- Date and Author Footer ---
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
    date_frame = date_box.text_frame
    date_para = date_frame.paragraphs[0]
    date_para.text = f"{author} | {datetime.now().strftime('%B %d, %Y')}"
    date_para.font.size = Pt(14)
    date_para.font.color.rgb = COLOR_LIGHT
    date_para.alignment = PP_ALIGN.CENTER
    
    return prs


def add_chart_slide(
    prs: Presentation,
    title: str,
    chart_path: str,
    description: str = None
) -> None:
    """
    Add a slide with a chart image.
    
    Creates a new slide with a title, embedded chart image, and optional
    description text at the bottom.
    
    Args:
        prs (Presentation): Presentation object to add slide to.
        title (str): Slide title.
        chart_path (str): Path to the chart image file.
        description (str, optional): Description text below the chart.
        
    Returns:
        None
        
    Example:
        >>> add_chart_slide(prs, "Revenue Trends", "charts/revenue.png", "Monthly comparison")
    """
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # --- Slide Title ---
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.75))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRIMARY
    
    # --- Chart Image ---
    # Centered with appropriate sizing
    slide.shapes.add_picture(chart_path, Inches(1.5), Inches(1.2), width=Inches(10))
    
    # --- Description Footer ---
    if description:
        desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.75))
        desc_frame = desc_box.text_frame
        desc_para = desc_frame.paragraphs[0]
        desc_para.text = description
        desc_para.font.size = Pt(14)
        desc_para.font.color.rgb = COLOR_SECONDARY
        desc_para.alignment = PP_ALIGN.CENTER


def add_summary_slide(
    prs: Presentation,
    title: str,
    summary_data: dict
) -> None:
    """
    Add a slide with summary statistics.
    
    Creates a slide displaying key metrics and statistics in a structured
    format. Supports nested dictionaries for grouped metrics.
    
    Args:
        prs (Presentation): Presentation object to add slide to.
        title (str): Slide title.
        summary_data (dict): Dictionary of metrics to display.
            Can contain nested dicts for grouped display.
            
    Returns:
        None
        
    Example:
        >>> summary = {
        ...     'Total Records': 1000,
        ...     'Revenue Metrics': {'revenue': {'Total': '$500K', 'Average': '$50K'}}
        ... }
        >>> add_summary_slide(prs, "Executive Summary", summary)
    """
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # --- Slide Title ---
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.75))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRIMARY
    
    # --- Summary Content ---
    y_position = 1.3  # Starting vertical position
    
    for key, value in summary_data.items():
        if isinstance(value, dict):
            # Section header for nested data
            header_box = slide.shapes.add_textbox(Inches(1), Inches(y_position), Inches(11), Inches(0.4))
            header_frame = header_box.text_frame
            header_para = header_frame.paragraphs[0]
            header_para.text = key.replace('_', ' ').title()
            header_para.font.size = Pt(18)
            header_para.font.bold = True
            y_position += 0.5
            
            # Nested items
            for sub_key, sub_value in value.items():
                if isinstance(sub_value, dict):
                    for metric, metric_value in sub_value.items():
                        item_box = slide.shapes.add_textbox(Inches(1.5), Inches(y_position), Inches(10), Inches(0.35))
                        item_frame = item_box.text_frame
                        item_para = item_frame.paragraphs[0]
                        
                        # Format numeric values
                        formatted_value = f"${metric_value:,.2f}" if isinstance(metric_value, (int, float)) else str(metric_value)
                        item_para.text = f"• {sub_key} {metric}: {formatted_value}"
                        item_para.font.size = Pt(14)
                        y_position += 0.35
        else:
            # Simple key-value pair
            item_box = slide.shapes.add_textbox(Inches(1), Inches(y_position), Inches(11), Inches(0.4))
            item_frame = item_box.text_frame
            item_para = item_frame.paragraphs[0]
            item_para.text = f"• {key.replace('_', ' ').title()}: {value}"
            item_para.font.size = Pt(16)
            y_position += 0.45


def save_presentation(prs: Presentation, filename: str) -> str:
    """
    Save the presentation to file.
    
    Args:
        prs (Presentation): Presentation object to save.
        filename (str): Output filename (saved in OUTPUT_DIR).
        
    Returns:
        str: Full path to the saved presentation file.
        
    Example:
        >>> path = save_presentation(prs, "report_2024.pptx")
    """
    filepath = OUTPUT_DIR / filename
    prs.save(filepath)
    return str(filepath)


# =============================================================================
# Module Test
# =============================================================================
if __name__ == "__main__":
    # Test presentation creation
    prs = create_presentation(
        "Monthly Sales Report",
        "Q4 2024 Performance Overview",
        "Data Analytics Team"
    )
    
    # Add a test summary slide
    test_summary = {
        'Total Records': 12,
        'Columns': 5,
    }
    add_summary_slide(prs, "Data Summary", test_summary)
    
    path = save_presentation(prs, "test_report.pptx")
    print(f"Test presentation saved: {path}")
